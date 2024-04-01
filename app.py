from flask import Flask, render_template, request, send_file
from pptx import Presentation
import os
from pptx.util import Pt
from pptx.dml.color import RGBColor
import logging

app = Flask(__name__)

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

@app.route('/')
def index():
    return render_template('index.html')

def modify_presentation(input_pptx, font_size, line_spacing, red, green, blue):
    try:
        prs = Presentation(input_pptx)

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            # Ensure font size is within valid range and convert to 100th of a point
                            font_size = max(1, min(font_size, 400000))  # Clamp font size to valid range
                            run.font.size = Pt(font_size)
                            run.font.color.rgb = RGBColor(red, green, blue)
                    shape.text_frame.paragraphs[0].line_spacing = line_spacing
                
        # Save output file
        output_path = os.path.join(app.root_path, 'modified_presentation.pptx')
        prs.save(output_path)
        
        return output_path  # Return the path of the modified file
    
    except Exception as e:
        logger.error("Error modifying presentation: %s", str(e))
        raise

@app.route('/process', methods=['POST'])
def process():
    try:
        # Get form data
        font_size = int(request.form.get('font_size', 12))
        line_spacing = float(request.form.get('line_spacing', 1))
        red = int(request.form.get('red', 0))
        green = int(request.form.get('green', 0))
        blue = int(request.form.get('blue', 0))

        # Get uploaded file
        ppt_file = request.files.get('ppt_file')
        if not ppt_file:
            return "No file uploaded", 400

        ppt_filename = ppt_file.filename

        # Save uploaded file
        upload_dir = 'uploads'
        if not os.path.exists(upload_dir):
            os.makedirs(upload_dir)
        ppt_path = os.path.join(upload_dir, ppt_filename)
        ppt_file.save(ppt_path)

        # Process PowerPoint file
        output_path = modify_presentation(ppt_path, font_size, line_spacing, red, green, blue)

        # Provide download link for modified file
        return send_file(output_path, as_attachment=True)
    except Exception as e:
        logger.error("Error processing file: %s", str(e))
        return "Error processing file", 500

if __name__ == '__main__':
    app.run(debug=True)
